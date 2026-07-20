"""Cargadores para Word, PDF, Excel y PowerPoint que devuelven documentos LangChain."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from docx import Document as WordDocument
from langchain_core.documents import Document
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".docx", ".pdf", ".xlsx", ".xls", ".pptx"}


class UnsupportedDocumentError(ValueError):
    """Indica que se solicitó un tipo de archivo que el proyecto no admite."""


def load_directory(directory: Path) -> tuple[Document, ...]:
    """Carga todos los documentos compatibles de una carpeta, ordenados por nombre."""

    documents: list[Document] = []
    for path in sorted(directory.iterdir(), key=lambda item: item.name.casefold()):
        if path.is_file() and path.suffix.casefold() in SUPPORTED_EXTENSIONS:
            documents.extend(load_file(path))
    return tuple(documents)


def load_file(path: Path) -> tuple[Document, ...]:
    """Lee un archivo y devuelve una o más unidades de contenido de LangChain."""

    extension = path.suffix.casefold()
    if extension == ".docx":
        return _load_word(path)
    if extension == ".pdf":
        return _load_pdf(path)
    if extension in {".xlsx", ".xls"}:
        return _load_excel(path)
    if extension == ".pptx":
        return _load_powerpoint(path)
    raise UnsupportedDocumentError(f"Formato no compatible: {path.name}")


def _metadata(path: Path, document_type: str, **extra: object) -> dict[str, object]:
    """Construye metadatos consistentes para trazabilidad de cada fragmento."""

    return {"source": str(path), "file_name": path.name, "file_type": document_type, **extra}


def _load_word(path: Path) -> tuple[Document, ...]:
    """Extrae párrafos y tablas de un archivo Word mediante python-docx."""

    word_file = WordDocument(path)
    parts = [paragraph.text.strip() for paragraph in word_file.paragraphs if paragraph.text.strip()]
    for table in word_file.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                parts.append(" | ".join(values))
    return (Document(page_content="\n".join(parts), metadata=_metadata(path, "word")),)


def _load_pdf(path: Path) -> tuple[Document, ...]:
    """Extrae cada página con texto de un PDF mediante pypdf."""

    reader = PdfReader(path)
    return tuple(
        Document(
            page_content=text,
            metadata=_metadata(path, "pdf", page=index + 1),
        )
        for index, page in enumerate(reader.pages)
        if (text := page.extract_text()).strip()
    )


def _load_excel(path: Path) -> tuple[Document, ...]:
    """Lee cada hoja de Excel con pandas y la conserva como texto tabular."""

    workbook = pd.read_excel(path, sheet_name=None)
    return tuple(
        Document(
            page_content=frame.to_csv(index=False),
            metadata=_metadata(path, "excel", sheet=sheet_name),
        )
        for sheet_name, frame in workbook.items()
        if not frame.empty
    )


def _load_powerpoint(path: Path) -> tuple[Document, ...]:
    """Extrae el texto visible de cada diapositiva de PowerPoint."""

    presentation = Presentation(path)
    documents: list[Document] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text:
            documents.append(
                Document(text, metadata=_metadata(path, "powerpoint", slide=index))
            )
    return tuple(documents)
